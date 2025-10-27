# app/presentation.py
from fastapi import APIRouter, HTTPException
from pydantic import BaseModel, validator
from typing import List, Optional
from io import BytesIO
from datetime import datetime
import hashlib
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from google.cloud import storage

router = APIRouter(prefix="/v1")

class GeneratePresentationRequest(BaseModel):
    collection: str
    title: Optional[str] = None
    date_of_birth: Optional[str] = None  # YYYY-MM-DD
    date_of_death: Optional[str] = None  # YYYY-MM-DD
    photos: List[str]
    output_bucket: str
    output_filename: str

    @validator("photos")
    def photos_not_empty(cls, v):
        if not isinstance(v, list) or len(v) == 0:
            raise ValueError("photos mag niet leeg zijn")
        return v

def _fmt_date(s: Optional[str]) -> Optional[str]:
    if not s: return None
    try:
        dt = datetime.strptime(s, "%Y-%m-%d")
        months = ["januari","februari","maart","april","mei","juni",
                  "juli","augustus","september","oktober","november","december"]
        return f"{dt.day} {months[dt.month-1]} {dt.year}"
    except Exception:
        return s

def _title_slide(prs: Presentation, title: Optional[str], dob: Optional[str], dod: Optional[str]):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    if title:
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.5))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    dob_fmt, dod_fmt = _fmt_date(dob), _fmt_date(dod)
    if dob_fmt or dod_fmt:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(8.4), Inches(1.0))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        if dob_fmt and dod_fmt:
            p2.text = f"Geboren {dob_fmt} — Overleden {dod_fmt}"
        elif dod_fmt:
            p2.text = f"Overleden {dod_fmt}"
        else:
            p2.text = f"Geboren {dob_fmt}"
        p2.font.size = Pt(20)
        p2.alignment = PP_ALIGN.CENTER
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

def _photo_slide(prs: Presentation, img_bytes: bytes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    stream = BytesIO(img_bytes)
    slide.shapes.add_picture(stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

def _download(url: str, timeout=20) -> bytes:
    r = requests.get(url, timeout=timeout)
    if r.status_code != 200:
        raise HTTPException(502, f"Foto niet te downloaden ({r.status_code})")
    return r.content

def _upload_gcs(bucket_name: str, blob_path: str, data: bytes, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation") -> str:
    client = storage.Client()
    bucket = client.bucket(bucket_name)
    blob = bucket.blob(blob_path)
    blob.upload_from_string(data, content_type=content_type)
    return f"https://storage.googleapis.com/{bucket_name}/{blob_path}"

@router.post("/generate-presentation")
def generate_presentation(req: GeneratePresentationRequest):
    try:
        if "/" in req.output_filename or req.output_filename.strip() == "":
            raise HTTPException(400, "Ongeldige output_filename")
        if req.collection.strip() == "":
            raise HTTPException(400, "collection is verplicht")

        prs = Presentation()  # 16:9
        _title_slide(prs, req.title, req.date_of_birth, req.date_of_death)

        for url in req.photos:  # volgorde Base44 aanhouden (jouw keuze B)
            try:
                img = _download(url)
                _photo_slide(prs, img)
            except HTTPException:
                continue

        if len(prs.slides) <= 1:
            raise HTTPException(400, "Geen geldige foto’s gevonden om te plaatsen")

        buf = BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        buf.close()

        h12 = hashlib.sha256(data).hexdigest()[:12]
        blob_path = f"{req.collection}/{h12}_{req.output_filename}"
        url = _upload_gcs(req.output_bucket, blob_path, data)

        return {"download_url": url}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Presentatiegeneratie mislukt: {str(e)}")
