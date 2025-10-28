import re
import os
import logging
import sys
import traceback
from dotenv import load_dotenv
from fastapi.middleware.cors import CORSMiddleware
from pptx.enum.shapes import PP_PLACEHOLDER

os.environ["PYTHONUNBUFFERED"] = "1"
load_dotenv()

# ✅ Logging direct naar stderr
logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(sys.stderr)],
    force=True
)

# ✅ Uvicorn logs ook op debug zetten
for logger_name in ["uvicorn", "uvicorn.error"]:
    logging.getLogger(logger_name).setLevel(logging.DEBUG)

# ✅ Ongecatchte exceptions afdrukken
def handle_exception(exc_type, exc_value, exc_traceback):
    if issubclass(exc_type, KeyboardInterrupt):
        return
    print("\n🔥 UNCAUGHT EXCEPTION 🔥", file=sys.stderr)
    traceback.print_exception(exc_type, exc_value, exc_traceback, file=sys.stderr)

sys.excepthook = handle_exception

from fastapi import FastAPI, HTTPException, Header
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Optional
from google.cloud import storage
from scripts.maak_presentatie import maak_presentatie_automatisch
from pydantic import BaseModel, validator
import requests
from io import BytesIO
import hashlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

API_KEY = os.getenv("STREAMLIT_API_KEY")
BUCKET_NAME = os.getenv("BUCKET_TEMPLATES")

logging.debug(f"✅ gestart met BUCKET_TEMPLATES = {BUCKET_NAME}")
logging.debug(f"✅ API_KEY loaded? {'✅' if API_KEY else '❌'}")

app = FastAPI(title="Uitvaart Presentatie API ✅")

# ✅ CORS toestaan (nodig voor Streamlit of externe clients)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ✅ BESTAANDE CLASS VOOR STREAMLIT ENDPOINT
class GenRequest(BaseModel):
    naam: str
    sjabloon: str
    fotos: List[str] = []
    geboortedatum: Optional[str] = None
    overlijdensdatum: Optional[str] = None

# ✅ NIEUWE CLASS VOOR BASE44 PRESENTATIE
class GeneratePresentationRequest(BaseModel):
    collection: str
    title: Optional[str] = None
    date_of_birth: Optional[str] = None
    date_of_death: Optional[str] = None
    photos: List[str]
    output_bucket: str
    output_filename: str
    template_file: Optional[str] = None

    @validator("photos")
    def photos_not_empty(cls, v):
        if not isinstance(v, list) or len(v) == 0:
            raise ValueError("Minimaal één foto verplicht")
        return v

@app.get("/")
def root():
    return {"status": "✅ API actief", "service": "Presentatie generator"}

def _fmt_date(s: Optional[str]) -> Optional[str]:
    if not s:
        return None
    try:
        dt = datetime.strptime(s, "%Y-%m-%d")
        months = ["januari","februari","maart","april","mei","juni",
                  "juli","augustus","september","oktober","november","december"]
        return f"{dt.day} {months[dt.month-1]} {dt.year}"
    except Exception:
        return s

@app.post("/v1/generate-presentation")
def generate_presentation(req: GeneratePresentationRequest):
    logging.debug(f"🚀 Base44 generate-presentation req: {req}")

    template_file = req.template_file or "SjabloonRustig.pptx"

    client = storage.Client()
    sjabloon_bucket = client.bucket("warmeuitvaartassistent-sjablonen")
    sjabloon_blob = sjabloon_bucket.blob(f"sjablonen/{template_file}")
    local_template = f"/tmp/sjablonen/{template_file}"
    os.makedirs("/tmp/sjablonen", exist_ok=True)

    if os.path.exists(local_template):
        os.remove(local_template)

    sjabloon_blob.reload()
    sjabloon_blob.download_to_filename(local_template)

    dob_fmt = _fmt_date(req.date_of_birth)
    dod_fmt = _fmt_date(req.date_of_death)

    if dob_fmt and dod_fmt:
        titel_datums = f"* {dob_fmt} – † {dod_fmt}"
    elif dob_fmt:
        titel_datums = f"* {dob_fmt}"
    elif dod_fmt:
        titel_datums = f"† {dod_fmt}"
    else:
        titel_datums = None

    try:
        logging.debug("🎬 PPT genereren gestart met sjabloon...")

        prs = Presentation(local_template)
        fotos = list(req.photos)
        foto_index = 0

        for slide in prs.slides:
            image_shapes = []
            for sh in slide.shapes:
                try:
                    if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                        image_shapes.append(sh)
                except Exception:
                    continue

            def _ph_order_key(sh):
                name = getattr(sh, "name", "") or ""
                m = re.search(r"foto\s*0*(\d+)", name, re.IGNORECASE)
                if m:
                    return (0, int(m.group(1)))
                return (1, int(sh.top), int(sh.left))

            image_shapes.sort(key=_ph_order_key)

            if not image_shapes:
                continue

            for placeholder in image_shapes:
                if foto_index >= len(fotos):
                    foto_index = 0
                try:
                    img_data = requests.get(fotos[foto_index]).content
                    placeholder.insert_picture(BytesIO(img_data))
                    foto_index += 1
                except Exception as e:
                    logging.error(f"❌ Foto kon niet worden geplaatst: {e}")
                    continue

        buf = BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        buf.close()

    except Exception as e:
        logging.exception("❌ Fout tijdens presentatie generatie")
        raise HTTPException(status_code=500, detail=str(e))

    bucket = client.bucket(req.output_bucket)

    import uuid
    unique_id = uuid.uuid4().hex[:8]
    blob_path = f"{req.collection}/{unique_id}_{req.output_filename}"

    blob = bucket.blob(blob_path)
    blob.upload_from_string(
        data,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    url = f"https://storage.googleapis.com/{req.output_bucket}/{blob_path}"
    logging.debug(f"✅ Downloadlink: {url}")

    return {"download_url": url}