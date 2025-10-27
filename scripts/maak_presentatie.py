# scripts/maak_presentatie.py
# -*- coding: utf-8 -*-
"""
Warme Uitvaartassistent — Automatische PowerPoint-bouwer met placeholder-herkenning
"""

import os
import re
import io
import zipfile
import tempfile
import shutil
import requests

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu
from PIL import Image


# ---------------------------
# Helpers voor bestanden/foto's
# ---------------------------

def verzamel_fotobestanden(uploadmap: str) -> list[str]:
    """Zoek alle .jpg/.jpeg/.png in een map (recursief)."""
    fotopaden = []
    for root, _, files in os.walk(uploadmap):
        for f in files:
            if f.lower().endswith((".jpg", ".jpeg", ".png")):
                fotopaden.append(os.path.join(root, f))
    return fotopaden


def download_base44_fotos(foto_urls: list[str], tmp_dir: str) -> list[str]:
    """Download Base44-foto's met retries en foutafhandeling."""
    import time
    paden = []

    for i, url in enumerate(foto_urls, start=1):
        success = False
        for attempt in range(3):
            try:
                r = requests.get(url, timeout=15)
                if r.status_code == 200:
                    ext = ".jpg"
                    ct = r.headers.get("Content-Type", "")
                    if "png" in ct.lower():
                        ext = ".png"
                    pad = os.path.join(tmp_dir, f"base44_foto_{i}{ext}")
                    with open(pad, "wb") as f:
                        f.write(r.content)
                    paden.append(pad)
                    success = True
                    break
            except Exception:
                pass
            time.sleep(1)
        if not success:
            print(f"⚠️ Kon foto {i} niet downloaden na 3 pogingen.")
    return paden


# ---------------------------
# Beeldverwerking
# ---------------------------

def _compute_contain_size(img_w: int, img_h: int, box_w_emu: Emu, box_h_emu: Emu) -> tuple[int, int]:
    """Bereken contain (fit) afmetingen."""
    box_w = int(box_w_emu)
    box_h = int(box_h_emu)
    r = min(box_w / img_w, box_h / img_h)
    return max(1, int(img_w * r)), max(1, int(img_h * r))


def _crop_to_ratio(img: Image.Image, target_w_emu: Emu, target_h_emu: Emu) -> Image.Image:
    """Crop de afbeelding zodat hij de shape volledig vult (cover)."""
    img_w, img_h = img.size
    target_ratio = float(target_w_emu) / float(target_h_emu)
    img_ratio = img_w / img_h

    if img_ratio > target_ratio:
        new_w = int(target_ratio * img_h)
        x0 = (img_w - new_w) // 2
        return img.crop((x0, 0, x0 + new_w, img_h))
    else:
        new_h = int(img_w / target_ratio)
        y0 = (img_h - new_h) // 2
        return img.crop((0, y0, img_w, y0 + new_h))


# ---------------------------
# Placeholder-detectie en vervanging
# ---------------------------

_PLACEHOLDER_RE = re.compile(r"^foto[_\-]?(\d+)$", re.IGNORECASE)


def _iter_all_shapes_recursive(container):
    """Itereer alle shapes, inclusief groepen."""
    for sh in container.shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for inner in _iter_all_shapes_recursive(sh):
                yield inner


def _collect_named_placeholders(prs: Presentation) -> list[tuple[int, object]]:
    """Zoek alle shapes met naam foto_x."""
    matches = []
    for slide in prs.slides:
        for sh in _iter_all_shapes_recursive(slide):
            nm = getattr(sh, "name", None)
            if not nm:
                continue
            m = _PLACEHOLDER_RE.match(nm)
            if m:
                idx = int(m.group(1))
                matches.append((idx, sh))
    matches.sort(key=lambda t: t[0])
    return matches


def _replace_shape_with_picture(slide, shape, image_path: str, ratio_mode: str = "cover"):
    """Vervang de visuele inhoud van een shape door een foto."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height

    if getattr(shape, "is_placeholder", False):
        try:
            if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                with Image.open(image_path) as im:
                    cropped = _crop_to_ratio(im, width, height)
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tf:
                        cropped.save(tf.name)
                        tmp = tf.name
                shape.insert_picture(tmp)
                os.remove(tmp)
                return
        except Exception:
            pass

    try:
        shape._element.getparent().remove(shape._element)
    except Exception:
        pass

    with Image.open(image_path) as im:
        cropped = _crop_to_ratio(im, width, height)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tf:
            cropped.save(tf.name)
            tmp = tf.name
    slide.shapes.add_picture(tmp, left, top, width=width, height=height)
    os.remove(tmp)


def vervang_placeholder_fotos(prs: Presentation, fotopaden: list[str], ratio_mode: str = "cover", repeat_if_insufficient: bool = True) -> int:
    """Vervang alle placeholders in het sjabloon."""
    placeholders = _collect_named_placeholders(prs)
    if not placeholders:
        print("Geen placeholders met naam foto_x gevonden in sjabloon.")
        return 0

    if not fotopaden:
        print("Geen fotopaden aangeleverd voor vervanging.")
        return 0

    totaal_fotos = len(fotopaden)
    vervangen = 0

    for i, (_, shape) in enumerate(placeholders):
        foto_pad = fotopaden[i % totaal_fotos] if repeat_if_insufficient else fotopaden[i] if i < totaal_fotos else None
        if not foto_pad:
            continue

        slide = None
        for s in prs.slides:
            if shape in s.shapes:
                slide = s
                break

        if slide:
            _replace_shape_with_picture(slide, shape, foto_pad, ratio_mode)
            vervangen += 1

    print(f"In totaal {vervangen} placeholders vervangen.")
    return vervangen


# ---------------------------
# Titel-dia
# ---------------------------

def zet_titel_dia(prs: Presentation, naam: str, datums: str | None = None, bijzin: str | None = None):
    """Vul titel en ondertitel op de eerste dia."""
    if not prs.slides:
        return
    slide = prs.slides[0]
    title_shape = None
    subtitle_shape = None

    for sh in slide.shapes:
        if getattr(sh, "is_placeholder", False):
            try:
                if sh.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    title_shape = sh
                elif sh.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                    subtitle_shape = sh
            except Exception:
                pass

    if title_shape and hasattr(title_shape, "text_frame"):
        title_shape.text_frame.clear()
        title_shape.text_frame.text = naam

    subtitle_lines = []
    if datums:
        subtitle_lines.append(datums)
    if bijzin:
        subtitle_lines.append(bijzin)

    if subtitle_shape and hasattr(subtitle_shape, "text_frame"):
        subtitle_shape.text_frame.clear()
        subtitle_shape.text_frame.text = "\n".join(subtitle_lines)


# ---------------------------
# Hoofdfunctie
# ---------------------------

def maak_presentatie_automatisch(
    sjabloon_pad: str,
    base44_foto_urls: list[str] | None = None,
    upload_bestanden: list[str] | None = None,
    uitvoer_pad: str = "uitvaart_presentatie_resultaat.pptx",
    ratio_mode: str = "cover",
    titel_naam: str | None = None,
    titel_datums: str | None = None,
    titel_bijzin: str | None = None,
    repeat_if_insufficient: bool = True
) -> str:
    """Bouw de presentatie en retourneer het pad naar het .pptx-bestand."""
    if not os.path.exists(sjabloon_pad):
        raise FileNotFoundError(f"Sjabloon niet gevonden: {sjabloon_pad}")

    if not base44_foto_urls and not upload_bestanden:
        raise ValueError("Geen invoer: geef base44_foto_urls of upload_bestanden op.")

    tmp_dir = tempfile.mkdtemp(prefix="presentatie_")
    print(f"Tijdelijke map aangemaakt: {tmp_dir}")

    fotopaden: list[str] = []

    try:
        if base44_foto_urls:
            fotopaden = download_base44_fotos(base44_foto_urls, tmp_dir)

        if upload_bestanden and not fotopaden:
            zip_bestanden = [f for f in upload_bestanden if f.lower().endswith(".zip")]
            if zip_bestanden:
                for zip_pad in zip_bestanden:
                    with zipfile.ZipFile(zip_pad, "r") as zip_ref:
                        zip_ref.extractall(tmp_dir)
                fotopaden = verzamel_fotobestanden(tmp_dir)
            else:
                for f in upload_bestanden:
                    if f.lower().endswith((".jpg", ".jpeg", ".png")):
                        doelpad = os.path.join(tmp_dir, os.path.basename(f))
                        shutil.copy(f, doelpad)
                        fotopaden.append(doelpad)

        if not fotopaden:
            raise ValueError("Geen geldige foto's gevonden om te verwerken.")

        prs = Presentation(sjabloon_pad)
        if titel_naam:
            zet_titel_dia(prs, titel_naam, titel_datums, titel_bijzin)

        vervang_placeholder_fotos(prs, fotopaden, ratio_mode=ratio_mode, repeat_if_insufficient=repeat_if_insufficient)

        OUTPUT_DIR = "/app/output"
        os.makedirs(OUTPUT_DIR, exist_ok=True)

        output_path = os.path.join(OUTPUT_DIR, uitvoer_pad)
        prs.save(output_path)

        return output_path


    except Exception as e:
        print(f"❌ Fout bij genereren of opslaan van de presentatie: {e}")
        raise

    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        print("🧹 Tijdelijke bestanden verwijderd.")
